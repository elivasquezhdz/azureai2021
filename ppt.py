from pptx import Presentation
import sys
import os
from pptx.util import Inches

uuid = sys.argv[1]

fldr = os.path.join("results",uuid)
imgs = [x for x in os.listdir(fldr) if x[-3:]=="png"]
prs = Presentation()

for i,img in enumerate(imgs):
    blank_slide_layout = prs.slide_layouts[i]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = top = Inches(1)
    height = Inches(5.5)
    
    

    pic = slide.shapes.add_picture(os.path.join(fldr,img), left, top,height=height)

prs.save(os.path.join("static",'{}.pptx'.format(uuid)))